#!/usr/bin/env python3
"""
PRODUCTION Re-indexing - Index ALL real documents from Azure Storage directly
"""

import os
import sys
import asyncio
import io
from dotenv import load_dotenv
from azure.storage.blob import BlobServiceClient
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
import re
from datetime import datetime
import time
import PyPDF2
import docx
from openai import AsyncAzureOpenAI
from azure.core.exceptions import ServiceResponseError
try:
    import openpyxl
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

try:
    from pptx import Presentation
    POWERPOINT_SUPPORT = True
except ImportError:
    POWERPOINT_SUPPORT = False

try:
    import olefile
    LEGACY_OFFICE_SUPPORT = True
except ImportError:
    LEGACY_OFFICE_SUPPORT = False

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))


def extract_pdf_content(blob_data: bytes) -> str:
    """Extract text content from PDF blob data with enhanced error reporting."""
    try:
        pdf_stream = io.BytesIO(blob_data)
        pdf_reader = PyPDF2.PdfReader(pdf_stream)
        
        # Check if PDF is encrypted
        if pdf_reader.is_encrypted:
            print(f"  📄 PDF is encrypted - attempting to decrypt...")
            try:
                # Try empty password first (common case)
                if pdf_reader.decrypt(""):
                    print(f"  ✅ Successfully decrypted with empty password")
                else:
                    print(f"  ❌ PDF requires password - cannot extract content")
                    return f"Encrypted PDF file (password required for content extraction)"
            except Exception as decrypt_error:
                print(f"  ❌ Decryption failed: {decrypt_error}")
                return f"Encrypted PDF file (decryption failed: {str(decrypt_error)})"
        
        num_pages = len(pdf_reader.pages)
        print(f"  📄 Processing PDF with {num_pages} pages...")
        
        text_content = ""
        pages_with_text = 0
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
                    pages_with_text += 1
            except Exception as e:
                print(f"  Warning: Could not extract page {page_num + 1}: {e}")
                continue
        
        if pages_with_text == 0:
            print(f"  📄 No extractable text found in {num_pages} pages (possibly scanned images)")
            return f"PDF file with {num_pages} pages (no extractable text - may be scanned images)"
        else:
            print(f"  ✅ Extracted text from {pages_with_text}/{num_pages} pages")
        
        return text_content.strip()
        
    except Exception as e:
        print(f"  ❌ PDF extraction error: {e}")
        return f"PDF file (extraction failed: {str(e)[:100]})"


def extract_docx_content(blob_data: bytes) -> str:
    """Extract text content from DOCX blob data with fallback options."""
    try:
        # First, try standard DOCX extraction
        docx_stream = io.BytesIO(blob_data)
        doc = docx.Document(docx_stream)
        
        text_content = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content += paragraph.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content += " | ".join(row_text) + "\n"
        
        return text_content.strip()
        
    except Exception as e:
        print(f"  Error extracting DOCX with python-docx: {e}")
        
        # Fallback 1: Try to extract as ZIP and read document.xml directly
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            docx_stream = io.BytesIO(blob_data)
            with zipfile.ZipFile(docx_stream, 'r') as zip_file:
                # Read the main document XML
                if 'word/document.xml' in zip_file.namelist():
                    xml_content = zip_file.read('word/document.xml')
                    root = ET.fromstring(xml_content)
                    
                    # Extract text from XML (basic extraction)
                    text_parts = []
                    for elem in root.iter():
                        if elem.text:
                            text_parts.append(elem.text)
                    
                    extracted_text = ' '.join(text_parts).strip()
                    if extracted_text:
                        print(f"  ✅ Extracted via XML fallback ({len(extracted_text)} chars)")
                        return extracted_text
                        
        except Exception as fallback_e:
            print(f"  Fallback XML extraction also failed: {fallback_e}")
        
        # Fallback 2: Try legacy Office extraction for potentially misnamed files
        try:
            legacy_content = extract_legacy_office_content(blob_data, 'doc')
            if legacy_content and len(legacy_content.strip()) > 20:
                print(f"  ✅ Extracted as legacy format ({len(legacy_content)} chars)")
                return legacy_content
        except Exception as legacy_e:
            print(f"  Legacy extraction also failed: {legacy_e}")
        
        # Fallback 3: Basic text extraction
        try:
            text_content = extract_text_content(blob_data)
            if text_content and len(text_content.strip()) > 10:
                print(f"  ✅ Extracted as plain text ({len(text_content)} chars)")
                return text_content
        except Exception as text_e:
            print(f"  Text extraction also failed: {text_e}")
        
        print(f"  ❌ All extraction methods failed for DOCX file")
        return f"DOCX file (extraction failed - possibly corrupted: {str(e)})"


def extract_legacy_office_content(blob_data: bytes, file_type: str) -> str:
    """Extract text content from legacy Office files (.doc, .xls, .ppt)."""
    try:
        # For legacy files, try to extract whatever text we can
        # This is basic extraction - not perfect but better than nothing
        text = blob_data.decode('utf-8', errors='ignore')
        
        # Remove common binary junk and control characters
        import re
        # Remove null bytes and control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', ' ', text)
        # Remove multiple spaces
        text = re.sub(r'\s+', ' ', text)
        # Extract potential text content (words with reasonable length)
        words = text.split()
        filtered_words = [word for word in words if len(word) > 2 and word.isalnum()]
        
        if len(filtered_words) > 10:  # If we found reasonable text
            return ' '.join(filtered_words[:1000])  # Limit to first 1000 words
        else:
            return f"Legacy {file_type.upper()} file (limited text extraction available)"
            
    except Exception as e:
        print(f"  Error extracting legacy {file_type}: {e}")
        return f"Legacy {file_type.upper()} file (content extraction failed: {str(e)})"


def extract_text_content(blob_data: bytes) -> str:
    """Extract text content from plain text files with binary detection."""
    try:
        # Check if content is likely binary by looking for null bytes and control characters
        sample_size = min(1024, len(blob_data))
        sample = blob_data[:sample_size]
        
        # If more than 5% of the sample contains null bytes or excessive control chars, it's likely binary
        null_count = sample.count(b'\x00')
        control_count = sum(1 for b in sample if b < 32 and b not in [9, 10, 13])  # Allow tab, LF, CR
        
        if null_count > sample_size * 0.01 or control_count > sample_size * 0.05:
            # This looks like binary content, don't try to decode
            return ""
        
        # Try UTF-8 first
        text = blob_data.decode('utf-8').strip()
        
        # Additional check: if decoded text has too many weird characters, skip it
        printable_chars = sum(1 for c in text if c.isprintable() or c.isspace())
        if len(text) > 0 and printable_chars / len(text) < 0.7:
            return ""
            
        return text
        
    except UnicodeDecodeError:
        try:
            # Fallback to latin-1, but still check for quality
            text = blob_data.decode('latin-1', errors='ignore').strip()
            
            # Check if the decoded text makes sense (not just binary garbage)
            printable_chars = sum(1 for c in text if c.isprintable() or c.isspace())
            if len(text) > 0 and printable_chars / len(text) < 0.7:
                return ""
                
            return text
        except Exception as e:
            print(f"  Error extracting text: {e}")
            return ""


def extract_excel_content(blob_data: bytes) -> str:
    """Extract text content from Excel files."""
    if not EXCEL_SUPPORT:
        return "Excel file (openpyxl not available for content extraction)"
    
    try:
        excel_stream = io.BytesIO(blob_data)
        workbook = openpyxl.load_workbook(excel_stream, data_only=True)
        
        text_content = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"\n--- Sheet: {sheet_name} ---\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    text_content += " | ".join(row_text) + "\n"
        
        return text_content.strip()
    except Exception as e:
        print(f"  Error extracting Excel: {e}")
        return f"Excel file (content extraction failed: {str(e)})"


def extract_powerpoint_content(blob_data: bytes) -> str:
    """Extract text content from PowerPoint files."""
    if not POWERPOINT_SUPPORT:
        return "PowerPoint file (python-pptx not available for content extraction)"
    
    try:
        ppt_stream = io.BytesIO(blob_data)
        prs = Presentation(ppt_stream)
        
        text_content = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text + "\n"
        
        return text_content.strip()
    except Exception as e:
        print(f"  Error extracting PowerPoint: {e}")
        return f"PowerPoint file (content extraction failed: {str(e)})"


def extract_rtf_content(blob_data: bytes) -> str:
    """Extract text content from RTF files (basic extraction)."""
    try:
        # Simple RTF parsing - just extract text between RTF commands
        text = blob_data.decode('utf-8', errors='ignore')
        
        # Remove RTF control codes and formatting
        import re
        # Remove RTF header and control words
        text = re.sub(r'\\[a-z]+\d*\s?', ' ', text)
        text = re.sub(r'[{}]', '', text)
        text = re.sub(r'\s+', ' ', text)
        
        return text.strip()
    except Exception as e:
        print(f"  Error extracting RTF: {e}")
        return f"RTF file (content extraction failed: {str(e)})"


def extract_outlook_message_content(blob_data: bytes, file_type: str) -> str:
    """Extract text content from Outlook message files (.msg, .eml)."""
    try:
        if file_type == 'eml':
            # EML files are standard email format (RFC 2822)
            import email
            from email.parser import BytesParser
            
            parser = BytesParser()
            msg = parser.parsebytes(blob_data)
            
            # Extract basic email metadata and content
            subject = msg.get('Subject', 'No Subject')
            sender = msg.get('From', 'Unknown Sender')
            recipient = msg.get('To', 'Unknown Recipient')
            date = msg.get('Date', 'Unknown Date')
            
            # Extract email body
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body += part.get_content()
                    elif part.get_content_type() == "text/html":
                        # Basic HTML to text conversion
                        html_content = part.get_content()
                        import re
                        # Remove HTML tags
                        body += re.sub(r'<[^>]+>', '', html_content)
            else:
                body = msg.get_content()
            
            # Combine metadata and content
            email_content = f"""
Subject: {subject}
From: {sender}
To: {recipient}
Date: {date}

{body}
            """.strip()
            
            return email_content
            
        elif file_type == 'msg':
            # MSG files are Outlook's proprietary format
            # Basic extraction - try to decode as much text as possible
            text = blob_data.decode('utf-8', errors='ignore')
            
            # Look for common email patterns and extract readable text
            import re
            
            # Try to find subject, sender, recipient patterns
            email_parts = []
            
            # Extract any readable text (this is basic - MSG format is complex)
            readable_text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)
            readable_text = re.sub(r'\s+', ' ', readable_text)
            
            # Look for email-like patterns
            subject_match = re.search(r'Subject:?\s*([^\n\r]+)', readable_text, re.IGNORECASE)
            if subject_match:
                email_parts.append(f"Subject: {subject_match.group(1).strip()}")
            
            from_match = re.search(r'From:?\s*([^\n\r]+)', readable_text, re.IGNORECASE)
            if from_match:
                email_parts.append(f"From: {from_match.group(1).strip()}")
            
            to_match = re.search(r'To:?\s*([^\n\r]+)', readable_text, re.IGNORECASE)
            if to_match:
                email_parts.append(f"To: {to_match.group(1).strip()}")
            
            # Add any substantial text content
            if len(readable_text.strip()) > 50:
                email_parts.append(f"Content: {readable_text.strip()[:2000]}")  # Limit to avoid too much noise
            
            return "\n".join(email_parts) if email_parts else f"Outlook message file (basic extraction: {readable_text[:500]})"
        
        return f"Email message file (format: {file_type})"
        
    except Exception as e:
        print(f"  Error extracting {file_type.upper()} message: {e}")
        return f"{file_type.upper()} message file (content extraction failed: {str(e)})"


def should_skip_file(blob_name: str) -> bool:
    """Check if file should be skipped based on extension or path."""
    filename = blob_name.lower()
    file_path = blob_name.lower()
    
    # Skip files in Trash folders (deleted files)
    trash_folders = ['/trash/', '\\trash\\', '/recycle/', '\\recycle\\', '/.trash/', '\\.trash\\']
    if any(trash_folder in file_path for trash_folder in trash_folders):
        return True
    
    # Skip files in backup folders
    backup_folders = ['/backup/', '\\backup\\', '/backups/', '\\backups\\', '/bak/', '\\bak\\']
    if any(backup_folder in file_path for backup_folder in backup_folders):
        return True
    
    # Skip files with backup naming patterns
    backup_patterns = ['-backup', '_backup', '.backup', '-bak', '_bak', 'backup-', 'backup_']
    if any(pattern in filename for pattern in backup_patterns):
        return True
    
    # Skip media files that don't contain searchable text
    skip_extensions = [
        # Video files
        '.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm', '.mkv', '.m4v',
        # Audio files
        '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',
        # Image files
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.svg', '.webp', '.ico',
        # Database and backup files
        '.accdb', '.mdb', '.accde', '.laccdb', '.ldb', '.bak', '.backup', '.db', '.sqlite', '.sqlite3',
        # Other binary/non-text files
        '.zip', '.rar', '.7z', '.tar', '.gz',
        '.exe', '.dll', '.bin', '.iso',
        '.psd', '.ai', '.eps',
        # Version control and placeholder files
        '.keep'
    ]
    
    return any(filename.endswith(ext) for ext in skip_extensions)


def extract_document_content(blob_name: str, blob_data: bytes) -> str:
    """Extract content from document based on file extension."""
    filename = blob_name.lower()
    
    # PDF files
    if filename.endswith('.pdf'):
        return extract_pdf_content(blob_data)
    
    # Modern Office formats (Office 2007+)
    elif filename.endswith(('.docx', '.dotx')):  # Include Word templates
        return extract_docx_content(blob_data)
    elif filename.endswith(('.xlsx', '.xlsm', '.xltx')):  # Include Excel templates
        return extract_excel_content(blob_data)
    elif filename.endswith(('.pptx', '.pptm', '.potx')):  # Include PowerPoint templates
        return extract_powerpoint_content(blob_data)
    
    # Legacy Office formats (Office 97-2003)
    elif filename.endswith('.doc'):
        return extract_legacy_office_content(blob_data, 'doc')
    elif filename.endswith('.xls'):
        return extract_legacy_office_content(blob_data, 'xls')
    elif filename.endswith('.ppt'):
        return extract_legacy_office_content(blob_data, 'ppt')
    
    # Other document formats
    elif filename.endswith('.rtf'):
        return extract_rtf_content(blob_data)
    elif filename.endswith(('.odt', '.ods', '.odp')):  # OpenOffice/LibreOffice
        return extract_legacy_office_content(blob_data, 'opendocument')
    
    # Email message files
    elif filename.endswith('.eml'):
        return extract_outlook_message_content(blob_data, 'eml')
    elif filename.endswith('.msg'):
        return extract_outlook_message_content(blob_data, 'msg')
    
    # Text-based files (including .txt which you specifically asked about)
    elif filename.endswith(('.txt', '.md', '.readme', '.csv', '.json', '.xml', '.html', '.htm', 
                           '.log', '.cfg', '.ini', '.yml', '.yaml', '.sql', '.py', '.js', '.css')):
        return extract_text_content(blob_data)
    
    else:
        # Try to extract as text first, fallback to metadata description
        try:
            text_content = extract_text_content(blob_data)
            if text_content and len(text_content.strip()) > 10:
                return text_content
        except:
            pass
        
        # Fallback to metadata description for truly unsupported files
        return f"File: {os.path.basename(blob_name)} (content extraction not supported for this file type)"


async def generate_embeddings(openai_client: AsyncAzureOpenAI, text: str) -> list:
    """Generate embeddings for the text content."""
    try:
        # Be more conservative with token limits (roughly 1 token per 4 characters)
        # 8192 tokens max, so use ~6000 characters to be safe
        max_chars = 6000
        truncated_text = text[:max_chars]
        
        response = await openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=truncated_text
        )
        return response.data[0].embedding
    except Exception as e:
        print(f"  Warning: Could not generate embeddings: {e}")
        return []


# Load environment variables from .env file
load_dotenv()

async def production_reindex():
    """Re-index ALL production documents directly from Azure Storage."""
    print("🚨 PRODUCTION RE-INDEXING - Processing ALL real documents")
    print("=" * 80)
    
    # Get settings from environment variables (same as production)
    connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
    
    # Correctly parse the search service name from the endpoint
    search_endpoint = os.getenv("AZURE_SEARCH_SERVICE_ENDPOINT")
    search_service_name = ""
    if search_endpoint:
        match = re.search(r"https://(.*?)\.search\.windows\.net", search_endpoint)
        if match:
            search_service_name = match.group(1)

    search_key = os.getenv("AZURE_SEARCH_ADMIN_KEY") or os.getenv("AZURE_SEARCH_API_KEY")
    index_name = os.getenv("AZURE_SEARCH_INDEX_NAME", "dtce-documents-index")
    container_name = os.getenv("AZURE_STORAGE_CONTAINER_NAME", "dtce-documents")
    
    if not connection_string or not search_key or not search_service_name:
        print("❌ Missing Azure credentials in environment variables")
        print("Set AZURE_STORAGE_CONNECTION_STRING, AZURE_SEARCH_SERVICE_ENDPOINT, and AZURE_SEARCH_ADMIN_KEY")
        return
    
    # Initialize clients
    storage_client = BlobServiceClient.from_connection_string(connection_string)
    search_client = SearchClient(
        endpoint=search_endpoint,
        index_name=index_name,
        credential=AzureKeyCredential(search_key)
    )
    
    # Initialize OpenAI client for embeddings
    openai_client = AsyncAzureOpenAI(
        api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        api_version="2024-02-15-preview",
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
    )
    
    container_client = storage_client.get_container_client(container_name)
    
    # Get an iterator for blobs from production storage
    print(f"📋 Getting blob iterator from container '{container_name}'...")
    try:
        blob_iterator = container_client.list_blobs()
    except Exception as e:
        print(f"❌ Failed to get blob iterator: {e}")
        return

    # Convert iterator to list and sort to prioritize Projects folder
    print(f"📋 Loading and sorting blobs to prioritize Projects folder...")
    try:
        all_blobs = list(blob_iterator)
        print(f"📊 Found {len(all_blobs)} total blobs")
        
        # Sort blobs: Projects folder first, then everything else
        def sort_key(blob):
            blob_name = blob.name.lower()
            if blob_name.startswith('projects/'):
                return (0, blob.name)  # Projects first
            else:
                return (1, blob.name)  # Everything else after
        
        sorted_blobs = sorted(all_blobs, key=sort_key)
        print(f"✅ Sorted blobs - Projects folder will be processed first")
        
    except Exception as e:
        print(f"❌ Failed to sort blobs: {e}")
        return

    # Index ALL documents (Projects first, then others)
    print(f"🔥 Indexing production documents (Projects folder prioritized)...")
    success_count = 0
    error_count = 0
    skipped_count = 0
    total_count = 0
    
    current_folder = None
    for blob in sorted_blobs:
        total_count += 1
        
        # Track folder changes to show progress
        blob_folder = blob.name.split('/')[0] if '/' in blob.name else 'Root'
        if current_folder != blob_folder:
            current_folder = blob_folder
            print(f"\n📁 Now processing folder: {current_folder}")
            if current_folder.lower() == 'projects':
                print(f"🎯 PRIORITY: Processing Projects folder first for architect documents")
        
        try:
            print(f"[{total_count}] {blob.name}")
            
            # Skip media files and other non-text files
            if should_skip_file(blob.name):
                print(f"  ⏭️  Skipping media/binary file")
                skipped_count += 1
                continue
            
            # Get blob metadata with retry
            blob_client = storage_client.get_blob_client(container=container_name, blob=blob.name)
            
            metadata = {}
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    blob_properties = blob_client.get_blob_properties()
                    metadata = blob_properties.metadata or {}
                    break
                except (ServiceResponseError, ConnectionResetError) as e:
                    if attempt < max_retries - 1:
                        print(f"  ... network error getting metadata, retrying in 5s ({e})")
                        time.sleep(5)
                    else:
                        raise e

            # Check if document already exists in search index with correct content
            document_id = re.sub(r'[^a-zA-Z0-9_-]', '_', blob.name)
            document_id = re.sub(r'_+', '_', document_id).strip('_')
            
            try:
                existing_doc = search_client.get_document(key=document_id)
                existing_content = existing_doc.get('content', '')
                
                # Check if content is already good - if so, skip this document entirely
                if (existing_content and 
                    len(existing_content) > 100 and
                    not existing_content.startswith('Document:') and
                    not existing_content.startswith('File:') and
                    'Document:' not in existing_content[:200]):  # Check first 200 chars for placeholder patterns
                    print(f"  ✅ Skipping - already has good content ({len(existing_content)} chars)")
                    skipped_count += 1
                    continue
                else:
                    # Content is bad placeholder content - needs reprocessing
                    if existing_content.startswith('Document:') or existing_content.startswith('File:'):
                        print(f"  🔄 Reprocessing - has placeholder content")
                    elif len(existing_content) <= 100:
                        print(f"  🔄 Reprocessing - content too short ({len(existing_content)} chars)")
                    else:
                        print(f"  🔄 Reprocessing - suspicious content pattern")
                        
            except Exception:
                # Document doesn't exist or error retrieving it - proceed with indexing
                print(f"  ➕ New document - indexing for first time")

            # Extract project info from the blob path itself as a fallback
            folder_path = blob.name.rsplit('/', 1)[0] if '/' in blob.name else ''
            project_name = ""
            year = None
            
            if folder_path:
                path_parts = folder_path.split("/")
                # Logic to find project name and year from path
                if len(path_parts) > 1 and path_parts[0] == 'Projects':
                    if len(path_parts) > 2 and path_parts[1].isdigit(): # e.g. Projects/220/
                        project_name = path_parts[1]
                        if len(path_parts) > 3 and path_parts[2].isdigit() and len(path_parts[2]) == 4:
                            year = int(path_parts[2])
                    elif len(path_parts) > 2: # e.g. Projects/Some Project/
                        project_name = path_parts[1]

            # Download and extract REAL content from the document
            print(f"  📄 Downloading and extracting content...")
            try:
                blob_data = blob_client.download_blob().readall()
                content = extract_document_content(blob.name, blob_data)
                
                if not content or len(content.strip()) < 50:
                    print(f"  ⚠️  Minimal content extracted ({len(content)} chars)")
                    # Fallback to metadata if content extraction fails
                    filename = os.path.basename(blob.name)
                    content = f"Document: {filename}"
                    if folder_path:
                        content += f" | Path: {folder_path}"
                    if project_name:
                        content += f" | Project: {project_name}"
                else:
                    print(f"  ✅ Extracted {len(content)} characters of content")
                    # Show preview of the extracted content
                    preview = content[:300].replace('\n', ' ').strip()
                    if len(content) > 300:
                        preview += "..."
                    print(f"  👁️  Preview: {preview}")
                
            except Exception as e:
                print(f"  ❌ Failed to download/extract content: {e}")
                # Fallback to metadata
                filename = os.path.basename(blob.name)
                content = f"Document: {filename} | Path: {folder_path} | Project: {project_name}"

            # Generate embeddings for the content
            content_vector = await generate_embeddings(openai_client, content)

            # Create search document
            document_id = re.sub(r'[^a-zA-Z0-9_-]', '_', blob.name)
            document_id = re.sub(r'_+', '_', document_id).strip('_')
            
            filename = os.path.basename(blob.name)
            
            search_document = {
                "id": document_id,
                "blob_name": blob.name,
                "blob_url": blob_client.url,
                "filename": filename,
                "content_type": metadata.get("content_type", blob.content_settings.content_type or ""),
                "folder": folder_path,
                "size": blob.size or 0,
                "content": content,
                "content_vector": content_vector,  # Add the embeddings
                "last_modified": blob.last_modified.isoformat(),
                "created_date": blob.creation_time.isoformat() if blob.creation_time else blob.last_modified.isoformat(),
                "project_name": project_name,
                "year": year
            }
            
            # Upload to search index with retry
            for attempt in range(max_retries):
                try:
                    result = search_client.upload_documents([search_document])
                    if result[0].succeeded:
                        success_count += 1
                        if total_count % 100 == 0:  # Progress every 100 docs
                            print(f"  ✅ Progress: {success_count} indexed, {skipped_count} skipped, {total_count} total")
                        break  # Break on success
                    else:
                        if attempt < max_retries - 1:
                            print(f"  ... upload failed, retrying in 5s. Reason: {result[0].error_message}")
                            time.sleep(5)
                        else:
                            error_count += 1
                            print(f"  ❌ Failed to index after retries. Reason: {result[0].error_message}")
                except (ServiceResponseError, ConnectionResetError) as e:
                    if attempt < max_retries - 1:
                        print(f"  ... network error during upload, retrying in 5s ({e})")
                        time.sleep(5)
                    else:
                        raise e
                
        except Exception as e:
            error_count += 1
            print(f"  ❌ Error processing {blob.name}: {str(e)[:150]}")
    
    print(f"\n🎉 PRODUCTION RE-INDEXING COMPLETE!")
    print(f"✅ Successfully indexed: {success_count}")
    print(f"⏭️  Skipped (already current): {skipped_count}")
    print(f"❌ Errors: {error_count}")
    print(f"📊 Total documents processed: {total_count}")
    if total_count > 0:
        print(f"📈 Processing rate: {((success_count + skipped_count)/total_count*100):.1f}%")
        print(f"🔥 Actually processed: {success_count} new/updated documents")
    
    if success_count > 0:
        print(f"\n🤖 Your production bot should now find documents!")
    else:
        print(f"\n💥 No documents were indexed - check credentials and permissions")

if __name__ == "__main__":
    asyncio.run(production_reindex())
